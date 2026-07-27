"""Render a deck dict (from ai/presentation.py) into a .pptx byte stream.

deck = {title, subtitle, slides: [{title, bullets: [str], notes: str}]}
"""

from __future__ import annotations

import io
from typing import Any


def deck_to_pptx(deck: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = deck.get("title") or "Presentation"
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = deck.get("subtitle") or ""

    bullet_layout = prs.slide_layouts[1]  # Title and Content
    for s in deck.get("slides", []):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.get("title") or ""

        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = s.get("bullets") or [""]
        for i, bullet in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(bullet)
            para.level = 0
            para.font.size = Pt(18)

        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def safe_filename(name: str) -> str:
    keep = "".join(c if c.isalnum() or c in " -_" else "_" for c in (name or "deck"))
    return keep.strip().replace(" ", "_")[:60] or "deck"
